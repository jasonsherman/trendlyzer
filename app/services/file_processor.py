"""
Service for processing different file types.
"""
from typing import Optional
from PyPDF2 import PdfReader
import docx
import pandas as pd
from docx import Document
import logging

logger = logging.getLogger(__name__)

def process_file(filepath: str, file_extension: str) -> Optional[str]:
    """Process different file types and extract their content."""
    try:
        # 1. Handle text-based formats
        if file_extension in ['txt', 'csv', 'md', 'rtf']:
            try:
                with open(filepath, 'r', encoding='utf-8') as f:
                    return f.read()
            except UnicodeDecodeError:
                with open(filepath, 'r', encoding='latin-1') as f:
                    return f.read()

        # 2. Handle PDFs
        elif file_extension == 'pdf':
            reader = PdfReader(filepath)
            return " ".join([page.extract_text() for page in reader.pages if page.extract_text()])

        # 3. Handle DOCX
        elif file_extension == 'docx':
            doc = docx.Document(filepath)
            return "\n".join([para.text for para in doc.paragraphs])

        # 4. Handle DOC
        elif file_extension == 'doc':
            doc = Document(filepath)
            text = []
            for para in doc.paragraphs:
                if para.text.strip():
                    text.append(para.text)
            for table in doc.tables:
                for row in table.rows:
                    row_data = [cell.text for cell in row.cells]
                    text.append('\t'.join(row_data))
            return "\n".join(text)

        # 5. Handle XLS/XLSX
        elif file_extension in ['xls', 'xlsx']:
            df = pd.read_excel(filepath, sheet_name=None)
            return "\n".join(df[sheet].to_string(index=False) for sheet in df)

        # 6. Handle PPT/PPTX
        elif file_extension in ['ppt', 'pptx']:
            from pptx import Presentation
            prs = Presentation(filepath)
            content = ""
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        content += shape.text + "\n"
            return content

        return None
    except Exception as e:
        logger.error(f"Error processing file: {e}")
        return None

def allowed_file(filename: str, allowed_extensions: set) -> bool:
    """Check if the file extension is allowed."""
    return '.' in filename and filename.rsplit('.', 1)[1].lower() in allowed_extensions 