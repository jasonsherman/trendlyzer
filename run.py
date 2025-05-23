"""
Trendlyzer - A microservice for analyzing documents and generating analytical reports with visualizations.
"""

# Standard library imports
from dotenv import load_dotenv
import matplotlib.pyplot as plt
from sumy.summarizers.lex_rank import LexRankSummarizer
from sumy.nlp.tokenizers import Tokenizer
from sumy.parsers.plaintext import PlaintextParser
from wordcloud import WordCloud
from nltk.corpus import stopwords
import nltk
import os
import logging
from collections import Counter, defaultdict
from typing import List, Dict, Tuple, Set, Optional
from dataclasses import dataclass
from textblob import TextBlob
from flask_mail import Mail, Message
import re
from sklearn.feature_extraction.text import TfidfVectorizer
from extract_text import extract_text
from llm_report import generate_report_with_llm
import base64

# Third-party imports
from flask import Flask, request, render_template, redirect, url_for, session, jsonify, send_file
from werkzeug.utils import secure_filename
import fitz  # PyMuPDF
from fpdf import FPDF
from fpdf.enums import XPos, YPos
import spacy
import matplotlib
matplotlib.use('Agg')

# Allowed extensions (no dot)
ALLOWED_EXTENSIONS = {'pdf', 'doc', 'docx', 'xls',
                      'xlsx', 'txt', 'csv', 'md', 'rtf', 'ppt', 'pptx'}

load_dotenv()
# Configure logging
logging.basicConfig(
    level=logging.INFO,
    format='%(asctime)s - %(name)s - %(levelname)s - %(message)s'
)
logger = logging.getLogger(__name__)

# Constants and Configuration
UPLOAD_FOLDER = 'uploads'
REPORTS_FOLDER = 'static/reports'
os.makedirs(UPLOAD_FOLDER, exist_ok=True)
os.makedirs(REPORTS_FOLDER, exist_ok=True)

# Report Configuration
REPORT_CONFIG = {
    'font': {
        'name': 'DejaVu',
        'regular': 'static/fonts/DejaVuSans.ttf',
        'bold': 'static/fonts/DejaVuSans-Bold.ttf',
        'italic': 'static/fonts/DejaVuSans-Oblique.ttf'
    },
    'colors': {
        'primary': (44, 82, 145),  # Dark blue
        'secondary': (77, 109, 243),  # Light blue
        'success': (76, 175, 80),  # Green
        'warning': (255, 152, 0),  # Orange
    },
    'charts': {
        'default_size': (4, 3),
        'bar_color': '#4d6df3',
        'max_percentage': 100
    }
}

# Theme Mapping
THEME_MAPPING = {
    'Lead Capture': ['email', 'phone', 'contact', 'address'],
    'Customer Support': ['assistance', 'support', 'help', 'question'],
    'AI Trust': ['ai', 'agent', 'human', 'real', 'bot'],
    'Sales Inquiry': ['price', 'cost', 'service', 'buy', 'purchase'],
    'Appointment Booking': ['appointment', 'schedule', 'meeting', 'consultation'],
}

# Known Companies and Locations
COMPANY_NAMES = ['vengo', 'purrfect', 'vengana']
LOCATIONS = [
    'serbia', 'usa', 'london', 'canada',
    'new york', 'germany', 'paris'
]

# Ensure NLTK can find the user's nltk_data directory
user_nltk_data = os.path.expanduser(
    r"C:/Users/spinnr/AppData/Roaming/nltk_data")
if user_nltk_data not in nltk.data.path:
    nltk.data.path.append(user_nltk_data)
print("NLTK data paths:", nltk.data.path)

# Download stopwords if not already present
try:
    _ = stopwords.words('english')
except LookupError:
    nltk.download('stopwords')
    _ = stopwords.words('english')

# Download punkt tokenizer if not already present
try:
    nltk.data.find('tokenizers/punkt')
except LookupError:
    nltk.download('punkt', download_dir=user_nltk_data)

try:
    nltk.data.find('tokenizers/punkt_tab')
except LookupError:
    nltk.download('punkt_tab')

# Helper for word frequency


def get_word_frequencies(text: str, top_n: int = 10) -> list:
    words = [w.lower() for w in re.findall(r'\b\w+\b', text)]
    stop_words = set(stopwords.words('english'))
    filtered = [w for w in words if w not in stop_words and len(w) > 2]
    return Counter(filtered).most_common(top_n)


def get_summary(text: str, sentence_count: int = 3) -> str:
    parser = PlaintextParser.from_string(text, Tokenizer("english"))
    summarizer = LexRankSummarizer()
    summary = summarizer(parser.document, sentence_count)
    return " ".join(str(sentence) for sentence in summary)


def improved_theme_detection(text: str, theme_mapping: dict, nlp_model) -> dict:
    doc = nlp_model(text)
    tokens = [
        token for token in doc if not token.is_stop and not token.is_punct and token.has_vector]
    theme_counts = {theme: 0 for theme in theme_mapping}
    for theme, keywords in theme_mapping.items():
        for kw in keywords:
            kw_vec = nlp_model(kw)[0]
            if not kw_vec.has_vector:
                continue
            # Count theme if any token in doc is similar to the keyword
            if any(kw_vec.similarity(token) > 0.75 for token in tokens):
                theme_counts[theme] += 1
    return {k: v for k, v in theme_counts.items() if v > 0}


def generate_wordcloud(text: str, output_path: str):
    stop_words = set(stopwords.words('english'))
    wc = WordCloud(width=800, height=400, background_color='white',
                   stopwords=stop_words).generate(text)
    wc.to_file(output_path)


@dataclass
class ReportMetrics:
    """Data class to hold report metrics."""
    word_count: int
    line_count: int
    total_conversations: int
    email_conversion_rate: float = 0.0
    phone_conversion_rate: float = 0.0
    follow_up_rate: float = 0.0
    readiness_rate: float = 0.0
    trust_rate: float = 0.0
    lead_success_rate: float = 0.0
    average_sentiment_score: float = 0.0

    mode: str = "Normal Document"


class ReportGenerator:
    """Class responsible for generating PDF reports with charts."""

    def __init__(self, filename: str, company_name: str):
        """Initialize the report generator.

        Args:
            filename: Name of the input file
            company_name: Name of the company
        """
        self.filename = filename
        self.company_name = company_name
        self.pdf = FPDF()
        self._setup_fonts()

    def _setup_fonts(self):
        """Register fonts for the PDF."""
        self.pdf.add_font(
            REPORT_CONFIG['font']['name'],
            '',
            REPORT_CONFIG['font']['regular'],
            uni=True
        )
        self.pdf.add_font(
            REPORT_CONFIG['font']['name'],
            'B',
            REPORT_CONFIG['font']['bold'],
            uni=True
        )
        self.pdf.add_font(
            REPORT_CONFIG['font']['name'],
            'I',
            REPORT_CONFIG['font']['italic'],
            uni=True
        )

    def _add_header(self):
        """Add header section to the report."""
        self.pdf.add_page()
        logo_path = 'static/images/trendlyzer-report-logo.png'

        if os.path.exists(logo_path):
            self.pdf.image(logo_path, x=80, y=10, w=50)
            self.pdf.ln(30)
        else:
            self.pdf.ln(20)

        self.pdf.set_font(REPORT_CONFIG['font']['name'], "B", 16)
        self.pdf.cell(
            0, 7, "Trend Analyzer & Insights Report",
            new_x=XPos.LMARGIN, new_y=YPos.NEXT, align='C'
        )

    def _create_chart(self, data: Dict[str, float], title: str, filename: str) -> str:
        """Create a bar chart and save it.

        Args:
            data: Dictionary of data points
            title: Chart title
            filename: Output filename

        Returns:
            str: Path to the saved chart
        """
        plt.figure(figsize=REPORT_CONFIG['charts']['default_size'])
        plt.bar(data.keys(), data.values(),
                color=REPORT_CONFIG['charts']['bar_color'])
        plt.title(title)
        plt.ylabel('Percentage (%)')
        plt.ylim(0, REPORT_CONFIG['charts']['max_percentage'])
        plt.xticks(rotation=20, ha='right')
        plt.tight_layout()

        chart_path = os.path.join(REPORTS_FOLDER, filename)
        os.makedirs(os.path.dirname(chart_path), exist_ok=True)
        plt.savefig(chart_path)
        plt.close()

        return chart_path

    def _add_section(self, title: str, content: str):
        """Add a section to the report.

        Args:
            title: Section title
            content: Section content
        """
        self.pdf.set_font(REPORT_CONFIG['font']['name'], "B", 14)
        self.pdf.set_text_color(*REPORT_CONFIG['colors']['primary'])
        self.pdf.cell(0, 7, title, new_x=XPos.LMARGIN, new_y=YPos.NEXT)
        self.pdf.set_text_color(0, 0, 0)
        self.pdf.set_font(REPORT_CONFIG['font']['name'], "", 12)
        self.pdf.multi_cell(0, 7, content)
        self.pdf.ln(8)

    def _add_metrics_chart(self, metrics: ReportMetrics):
        """Add metrics chart to the report.

        Args:
            metrics: Report metrics
        """
        chart_data = {
            'Email Leads': metrics.email_conversion_rate,
            'Phone Numbers': metrics.phone_conversion_rate,
            'Follow-Ups': metrics.follow_up_rate
        }

        chart_path = self._create_chart(
            chart_data,
            'Lead Capture Metrics',
            f"{os.path.basename(self.filename).replace('.txt', '')}_metrics.png"
        )

        self.pdf.image(chart_path, w=110)
        self.pdf.ln(10)

    def _add_keywords_section(self, top_keywords: List[Tuple[str, int]]):
        """Add keywords section to the report.

        Args:
            top_keywords: List of top keywords with counts
        """
        self._add_section("Top Keywords", "")
        for kw, count in top_keywords:
            self.pdf.cell(
                0, 7, f"- {kw} ({count} mentions)",
                new_x=XPos.LMARGIN, new_y=YPos.NEXT
            )

    def _add_themes_section(self, theme_counts: Dict[str, int]):
        """Add themes section to the report.

        Args:
            theme_counts: Dictionary of theme counts
        """
        self._add_section("Detected Themes", "")
        for theme, count in theme_counts.items():
            self.pdf.cell(
                0, 7, f"- {theme} ({count} mentions)",
                new_x=XPos.LMARGIN, new_y=YPos.NEXT
            )

    def _add_footer(self):
        """Add footer to the report."""
        self.pdf.set_y(self.pdf.get_y() + 5)
        self.pdf.set_font(REPORT_CONFIG['font']['name'], 'I', 8)
        self.pdf.cell(
            0, 10, 'Generated by Trendlyzer | VengoAI.com', 0, 1, 'C')

    def _add_wordcloud(self, text: str):
        wc_path = os.path.join(
            REPORTS_FOLDER, f"{os.path.basename(self.filename).replace('.txt', '')}_wordcloud.png")
        generate_wordcloud(text, wc_path)
        self.pdf.image(wc_path, w=110)
        self.pdf.ln(10)

    def _add_bar_chart(self, word_freqs: list):
        words, counts = zip(*word_freqs) if word_freqs else ([], [])
        plt.figure(figsize=(6, 3))
        plt.bar(words, counts, color=REPORT_CONFIG['charts']['bar_color'])
        plt.title('Top Words')
        plt.ylabel('Frequency')
        plt.xticks(rotation=20, ha='right')
        plt.tight_layout()
        chart_path = os.path.join(
            REPORTS_FOLDER, f"{os.path.basename(self.filename).replace('.txt', '')}_topwords.png")
        os.makedirs(os.path.dirname(chart_path), exist_ok=True)
        plt.savefig(chart_path)
        plt.close()
        self.pdf.image(chart_path, w=110)
        self.pdf.ln(10)

    def generate(
        self,
        mode: str,
        metrics: ReportMetrics,
        top_keywords: List[Tuple[str, int]],
        theme_counts: Dict[str, int],
        conversations: List[str],
        full_text: str
    ) -> Tuple[str, str]:
        try:
            self._add_header()
            if mode == "Conversational Document":
                overview = (
                    f"This report was created for {self.company_name} as a conversational document containing approximately "
                    f"{metrics.word_count:,} words and {metrics.line_count:,} lines. The analysis identified top themes like {', '.join(list(theme_counts.keys())[:3])}."
                )
                self._add_section("Overview", overview)
                highlights = (
                    f"- Total Conversations Analyzed: {metrics.total_conversations}\n"
                    f"- Email Leads Collected: {metrics.email_conversion_rate:.2f}%\n"
                    f"- Phone Numbers Collected: {metrics.phone_conversion_rate:.2f}%"
                )
                self._add_section("Key Highlights", highlights)
                # Lead Capture Rates Chart
                chart_data = {
                    'Email Leads': metrics.email_conversion_rate,
                    'Phone Numbers': metrics.phone_conversion_rate
                }
                chart_path = self._create_chart(
                    chart_data,
                    'Lead Capture Rates',
                    f"{os.path.basename(self.filename).replace('.txt', '')}_lead_capture.png"
                )
                self.pdf.image(chart_path, w=110)
                self.pdf.ln(10)
                # Top 3 Lead Capture Metrics
                lead_metrics = (
                    f"- {metrics.email_conversion_rate:.2f}% of customers provided an email after chatting\n"
                    f"- {metrics.phone_conversion_rate:.2f}% of customers provided a phone number\n"
                    f"- Over {metrics.follow_up_rate:.0f}% of all conversations led to actionable follow-ups"
                )
                self._add_section("Top 3 Lead Capture Metrics", lead_metrics)
                # Top Lead Capture Metrics Chart
                lead_chart_data = {'Email Provided': metrics.email_conversion_rate,
                                   'Phone Provided': metrics.phone_conversion_rate, 'Follow-Ups': metrics.follow_up_rate}
                lead_chart_path = self._create_chart(
                    lead_chart_data,
                    'Top Lead Capture Metrics',
                    f"{os.path.basename(self.filename).replace('.txt', '')}_top_lead_metrics.png"
                )
                self.pdf.image(lead_chart_path, w=110)
                self.pdf.ln(10)
                # Top Keywords and Conversation Drivers
                drivers_intro = (
                    "Analyzing word frequency provides insight into what customers care about most. "
                    "These keywords reveal user intents, buying readiness, and support needs."
                )
                self._add_section(
                    "Top Keywords and Conversation Drivers", drivers_intro)
                for kw, count in top_keywords:
                    self.pdf.cell(
                        0, 7, f"- {kw} ({count} mentions)", new_x=XPos.LMARGIN, new_y=YPos.NEXT)
                self.pdf.ln(8)
                # Trends & Impact on Businesses
                trends = (
                    "- Higher Lead Volume: Businesses using AI agents are collecting contact information in ~78% of conversations, significantly above the industry average.\n"
                    "- Customer Readiness: Nearly half of conversations involved direct requests to speak to a human or book services, showing strong buying intent.\n"
                    "- Trust Signals Matter: Many customers questioned whether the AI agent was real or human, highlighting the need for credibility-building in conversations."
                )
                self._add_section("Trends & Impact on Businesses", trends)
                # Visual Insights Chart (Business Trends Observed)
                vis_chart_data = {'Lead Capture Success': metrics.lead_success_rate,
                                  'Customer Readiness': metrics.readiness_rate, 'Trust Concerns': metrics.trust_rate}
                vis_chart_path = self._create_chart(
                    vis_chart_data,
                    'Business Trends Observed',
                    f"{os.path.basename(self.filename).replace('.txt', '')}_trends_impact.png"
                )
                self.pdf.image(vis_chart_path, w=110)
                self.pdf.ln(10)
                # Top Keywords
                self._add_section("Top Keywords", "")
                for kw, count in top_keywords:
                    self.pdf.cell(
                        0, 7, f"- {kw} ({count} mentions)", new_x=XPos.LMARGIN, new_y=YPos.NEXT)
                self.pdf.ln(8)
                # Detected Themes
                self._add_section("Detected Themes", "")
                for theme, count in theme_counts.items():
                    self.pdf.cell(
                        0, 7, f"- {theme} ({count} mentions)", new_x=XPos.LMARGIN, new_y=YPos.NEXT)
                self.pdf.ln(8)
                # Common Themes in Conversations
                common_themes = [
                    "Request Service Information: Detailed questions about services, pricing, timelines, and processes.",
                    "Initiate Sales Inquiries: Many customers express interest in purchasing or getting started.",
                    "Seek Customer Support: Some customers come with support-related questions, often post-sale.",
                    "Schedule Appointments or Demos: Agents were frequently asked how to book consultations or services.",
                    "Verify Business Credibility: Visitors often ask if the company is real, legitimate, or human-operated.",
                    "Location and Availability: Customers ask where the business operates or if service is available in their area."
                ]
                self._add_section("Common Themes in Conversations", "\n".join(
                    f"- {t}" for t in common_themes))
                # Sales Insights & Strategies
                sales_insights = [
                    "AI agents are effectively pre-qualifying leads by gathering needs, budget information, and location details.",
                    "Teams can intervene only when necessary, saving significant time and resources while maintaining engagement quality.",
                    "Businesses can leverage the data gathered during chats to personalize follow-up outreach and improve closing rates."
                ]
                self._add_section("Sales Insights & Strategies", "\n".join(
                    f"- {s}" for s in sales_insights))
                # Recommendations
                recommendations = [
                    "Add clear Call-to-Actions (CTAs) to AI flows (e.g., 'Would you like a team member to reach out?') to increase lead capture rates.",
                    "Personalize follow-up emails using data collected during AI conversations for higher engagement.",
                    "Refine agent scripts by focusing on top keywords and conversation intents to streamline answers and boost customer trust."
                ]
                self._add_section("Recommendations", "\n".join(
                    f"- {r}" for r in recommendations))
                self._add_footer()
            else:
                overview = (
                    f"This report was created for {self.company_name} as a business document containing approximately "
                    f"{metrics.word_count:,} words and {metrics.line_count:,} lines. The analysis identified top themes like {', '.join(list(theme_counts.keys())[:3])}."
                )
                self._add_section("Overview", overview)
                highlights = (
                    f"- Total Lines: {metrics.line_count}\n"
                    f"- Total Words: {metrics.word_count}"
                )
                self._add_section("Document Highlights", highlights)
                # Add summary
                summary = get_summary(full_text)
                self._add_section("Summary", summary)
                # Add word frequency bar chart
                word_freqs = get_word_frequencies(full_text, top_n=10)
                self._add_bar_chart(word_freqs)
                # Add word cloud
                self._add_wordcloud(full_text)
                # Top Keywords
                self._add_section("Top Keywords", "")
                for kw, count in top_keywords:
                    self.pdf.cell(
                        0, 7, f"- {kw} ({count} mentions)", new_x=XPos.LMARGIN, new_y=YPos.NEXT)
                self.pdf.ln(8)
                # Detected Themes
                self._add_section("Detected Themes", "")
                for theme, count in theme_counts.items():
                    self.pdf.cell(
                        0, 7, f"- {theme} ({count} mentions)", new_x=XPos.LMARGIN, new_y=YPos.NEXT)
                self.pdf.ln(8)
                # Recommendations
                recommendations = [
                    "Add clear Call-to-Actions (CTAs) to business flows to increase engagement.",
                    "Personalize follow-up emails using data collected during business interactions for higher engagement.",
                    "Refine business documents by focusing on top keywords and intents to streamline communication and boost trust."
                ]
                self._add_section("Recommendations", "\n".join(
                    f"- {r}" for r in recommendations))
                self._add_footer()
            if not os.path.exists(REPORTS_FOLDER):
                os.makedirs(REPORTS_FOLDER)
            report_filename = f"{os.path.basename(self.filename).replace('.txt', '')}_report.pdf"
            report_path = os.path.join(REPORTS_FOLDER, report_filename)
            self.pdf.output(report_path)
            web_report_path = f"/static/reports/{report_filename}"
            return web_report_path, overview
        except Exception as e:
            logger.error(f"Error generating report: {e}")
            raise


# Initialize Flask app
app = Flask(__name__)
app.config['UPLOAD_FOLDER'] = UPLOAD_FOLDER
app.secret_key = 'your_secret_key_here'

app.config['MAIL_SERVER'] = 'smtp.gmail.com'
app.config['MAIL_PORT'] = 587
app.config['MAIL_USE_TLS'] = True
app.config['MAIL_USERNAME'] = os.getenv("MAIL_USERNAME")
# your email password or app password
app.config['MAIL_PASSWORD'] = os.getenv("MAIL_PASSWORD")
app.config['MAIL_DEFAULT_SENDER'] = os.getenv("MAIL_USERNAME")

mail = Mail(app)
try:
    try:
        nlp = spacy.load("en_core_web_md")
        logger.info("Loaded spaCy model: en_core_web_md (with word vectors)")
    except Exception as e:
        logger.warning(
            "Could not load en_core_web_md, falling back to en_core_web_sm. Similarity may be less accurate.")
        nlp = spacy.load("en_core_web_sm")
        logger.info("Loaded spaCy model: en_core_web_sm (no word vectors)")
    logger.info("NLP models loaded successfully")
except Exception as e:
    logger.error(f"Error loading NLP models: {e}")
    raise


def extract_keywords_tfidf(text: str, top_n: int = 5) -> List[Tuple[str, float]]:
    """Enhanced TF-IDF keyword extraction with additional preprocessing."""
    # Initialize TF-IDF Vectorizer with more options
    vectorizer = TfidfVectorizer(
        stop_words='english',
        max_features=top_n,
        ngram_range=(1, 2),  # Allow both single words and 2-word phrases
        strip_accents='unicode',
        lowercase=True
    )

    try:
        # Clean and preprocess text
        text = ' '.join(text.split())  # Normalize whitespace

        # Fit and transform
        X = vectorizer.fit_transform([text])
        feature_names = vectorizer.get_feature_names_out()
        scores = X.toarray()[0]

        # Create and sort keywords
        keywords = list(zip(feature_names, scores))
        keywords.sort(key=lambda x: x[1], reverse=True)

        return keywords[:top_n]
    except Exception as e:
        logger.error(f"Error in TF-IDF keyword extraction: {e}")
        return []


def allowed_file(filename):
    # Get extension without dot, lowercased
    ext = os.path.splitext(filename)[1][1:].lower()
    print("Uploaded filename:", filename, "Extension:", ext)  # Debug print
    return '.' in filename and ext in ALLOWED_EXTENSIONS


def split_conversations(text: str) -> List[str]:
    """Split text into individual conversations. Each conversation starts with 'Agent:' at the start of a line."""
    # Split on every 'Agent:' at the start of a line, but keep the delimiter
    parts = re.split(r'(?=^Agent:)', text, flags=re.MULTILINE)
    # Remove any empty or whitespace-only conversations
    conversations = [p.strip() for p in parts if p.strip()]
    return conversations


def categorize_keyword(kw: str) -> str:
    """Categorize a keyword based on predefined lists.

    Args:
        kw: The keyword to categorize

    Returns:
        str: The category of the keyword
    """
    kw_lower = kw.lower()
    if kw_lower in COMPANY_NAMES:
        return 'company name'
    elif kw_lower in LOCATIONS:
        return 'location'
    return kw


def extract_text_from_pdf(file_path: str) -> str:
    """Extract text content from a PDF file.

    Args:
        file_path: Path to the PDF file

    Returns:
        str: Extracted text content
    """
    text = ""
    with fitz.open(file_path) as doc:
        for page in doc:
            text += page.get_text()
    return text


def create_conv_record(conv_id, user_name):
    return {
        "Conversation ID": conv_id,
        "User": user_name,
        "Email Captured": False,
        "Phone Captured": False,
        "Lead Capture Success": False,  # Will compute later
        "Follow‑up": False,
        "Customer Readiness": False,
        "Trust Concerns": False,
        "Sentiment Score": 0.0,
        "Message Count": 0
    }


@app.route('/')
def home() -> str:
    """Render the home page.

    Returns:
        str: Rendered HTML template
    """
    return render_template('index.html')


def send_email(report_path, company_name) -> str:
    try:
        recipient_email = os.getenv("RECEIVER_MAIL")
        msg = Message(
            subject=f"Trendlyzer Report for {company_name}",
            recipients=[recipient_email],
            body=f"A Trendlyzer report for {company_name} was just analyzed!"
        )
        file_path = report_path.lstrip('/')
        with open(file_path, 'rb') as fp:
            msg.attach(
                filename=os.path.basename(file_path),
                content_type='application/pdf',
                data=fp.read()
            )
        mail.send(msg)
        logger.info(f"Report sent to {recipient_email}")
        return f"EMAIL sent to {recipient_email} with {report_path}"
    except Exception as e:
        logger.error(f"Failed to send email: {e}")


@app.route('/upload', methods=['POST'])
def upload_file() -> str:
    """Handle file upload and process the uploaded file.

    Returns:
        str: Rendered template with results or error message
    """
    if 'file' not in request.files:
        return 'No file part in the request.'

    file = request.files['file']
    company_name = request.form.get(
        'company_name', 'Company Name not provided')

    if file.filename == '':
        return 'No selected file.'

    if file and allowed_file(file.filename):
        try:
            filename = secure_filename(file.filename)
            filepath = os.path.join(app.config['UPLOAD_FOLDER'], filename)
            os.makedirs(app.config['UPLOAD_FOLDER'], exist_ok=True)
            file.save(filepath)
            ext = os.path.splitext(filename)[1][1:].lower()

            text = extract_text(filepath, ext)
            if text is None:
                return f'File {file.filename} uploaded, but advanced analysis not yet available.'

            word_count = len(text.split())
            line_count = len(text.splitlines())

            lines = [line.rstrip("\n") for line in text.splitlines()]
            # Detect if it's a conversational document
            has_agent = any(re.match(r"Agent:", line) for line in lines)
            has_other_speaker = any(re.match(
                r"[^:]{1,40}:", line) and not line.startswith("Agent:") for line in lines)
            mode = "Conversational Document" if has_agent and has_other_speaker else "Normal Document"

            conversation_ids = []
            conversations = []
            current_conv_id = -1
            current_user = None
            current_conversation = ""
            all_keywords = []

            email_pattern = re.compile(
                r"[a-zA-Z0-9._%+-]+@[a-zA-Z0-9.-]+\.[a-zA-Z]{2,}")
            phone_pattern = re.compile(
                r"(?:\+?\d{1,3}[-.\s]?)?(?:\(?\d{2,4}\)?[-.\s]?)?\d{3}[-.\s]?\d{3,4}[-.\s]?\d{0,4}")
            followup_keywords = re.compile(
                r"\b(follow up|schedule|demo|call|reach out|appointment|book)\b", re.IGNORECASE)
            readiness_keywords = re.compile(
                r"\b(buy|purchase|ready|interested|go ahead|sign me up|subscribe|order|start|proceed)\b", re.IGNORECASE)
            trust_keywords = re.compile(
                r"\b(scam|fake|trust|secure|safety|safe|legit|fraud|privacy|data leak|security)\b", re.IGNORECASE)

            conv_data = []

            # Parse the file to identify conversations and collect stats
            for line in lines:
                if not line.strip() or ":" not in line:
                    continue  # skip empty and malformed lines
                speaker, message = [x.strip() for x in line.split(":", 1)]

                current_conversation += message

                if speaker == "Agent":
                    # If no conversation started yet, skip until user speaks
                    if current_conv_id < 0:
                        continue
                    # Follow‑up detection
                    if followup_keywords.search(message):
                        conv_data[current_conv_id]["Follow‑up"] = True

                    # Sentiment for agent message
                    conv_data[current_conv_id]["Sentiment Score"] += TextBlob(
                        message).sentiment.polarity
                    conv_data[current_conv_id]["Message Count"] += 1

                else:
                    if speaker != current_user:
                        current_user = speaker
                        current_conv_id += 1
                        conversation_ids.append(current_conv_id)
                        conversations.append(current_conversation)
                        conv_data.append(create_conv_record(
                            current_conv_id, current_user))

                        keywords = extract_keywords_tfidf(
                            current_conversation, top_n=5)
                        convo_keywords = [categorize_keyword(
                            kw) for kw, _ in keywords]

                        all_keywords.extend(convo_keywords)
                        current_conversation = ""

                    # Email / phone
                    if email_pattern.search(message):
                        conv_data[current_conv_id]["Email Captured"] = True
                    if phone_pattern.search(message):
                        nums = re.sub(
                            r"\D", "", phone_pattern.search(message).group())
                        if len(nums) >= 7:
                            conv_data[current_conv_id]["Phone Captured"] = True
                    # Readiness
                    if readiness_keywords.search(message):
                        conv_data[current_conv_id]["Customer Readiness"] = True

                    # Trust concerns
                    if trust_keywords.search(message):
                        conv_data[current_conv_id]["Trust Concerns"] = True

                    # Sentiment
                    conv_data[current_conv_id]["Sentiment Score"] += TextBlob(
                        message).sentiment.polarity
                    conv_data[current_conv_id]["Message Count"] += 1

            top_keywords = Counter(all_keywords).most_common(10)
            theme_counts = {}
            for kw, count in top_keywords:
                for theme, keywords in THEME_MAPPING.items():
                    if kw.lower() in keywords:
                        theme_counts[theme] = theme_counts.get(
                            theme, 0) + count

            # For business docs, use improved theme detection
            if mode == "Conversational Document":
                total_conversations = len(conversations)
                for d in conv_data:
                    d["Lead Capture Success"] = d["Email Captured"] or d["Phone Captured"]
                    # Average sentiment
                    if d["Message Count"]:
                        d["Sentiment Score"] = round(
                            d["Sentiment Score"] / d["Message Count"], 3)
                # Aggregate metrics
                lead_success_count = sum(
                    d["Lead Capture Success"] for d in conv_data)
                readiness_count = sum(d["Customer Readiness"]
                                      for d in conv_data)
                trust_count = sum(d["Trust Concerns"] for d in conv_data)

                lead_success_rate = (
                    lead_success_count / total_conversations * 100) if total_conversations else 0
                readiness_rate = (
                    readiness_count / total_conversations * 100) if total_conversations else 0
                trust_rate = (trust_count / total_conversations *
                              100) if total_conversations else 0

                email_leads = sum(d["Email Captured"] for d in conv_data)
                phone_leads = sum(d["Phone Captured"] for d in conv_data)
                followup_conv_count = sum(d["Follow‑up"] for d in conv_data)

                email_conversion_rate = (
                    email_leads / total_conversations) * 100 if total_conversations else 0
                phone_conversion_rate = (
                    phone_leads / total_conversations) * 100 if total_conversations else 0
                follow_up_rate = (
                    followup_conv_count / total_conversations * 100) if total_conversations else 0
                theme_counts_final = theme_counts
            else:
                email_conversion_rate = 0
                phone_conversion_rate = 0
                follow_up_rate = 0
                readiness_rate = 0
                lead_success_rate = 0
                trust_rate = 0
                total_conversations = 0
                theme_counts_final = improved_theme_detection(
                    text, THEME_MAPPING, nlp)

            metrics = ReportMetrics(
                word_count=word_count,
                line_count=line_count,
                total_conversations=total_conversations,
                email_conversion_rate=round(email_conversion_rate, 2),
                phone_conversion_rate=round(phone_conversion_rate, 2),
                follow_up_rate=round(follow_up_rate, 2),
                readiness_rate=round(readiness_rate, 2),
                trust_rate=round(trust_rate, 2),
                lead_success_rate=round(lead_success_rate, 2),
                average_sentiment_score=round(sum(
                    d["Sentiment Score"] for d in conv_data) / total_conversations, 3) if total_conversations else 0,
                mode=mode
            )

            report_generator = ReportGenerator(filename, company_name)
            report_path, overview = report_generator.generate(
                mode=mode,
                metrics=metrics,
                top_keywords=top_keywords,
                theme_counts=theme_counts_final,
                conversations=conversations,
                full_text=text
            )
            session['results'] = {
                'company_name': company_name,
                'overview': overview,
                'top_keywords': top_keywords,
                'theme_counts': theme_counts_final,
                'report_path': report_path
            }
            send_email(report_path, company_name)
            return redirect(url_for('results_page'))

        except Exception as e:
            logger.error(f"Error processing file: {e}")
            return f"An error occurred while processing the file: {str(e)}"

    return 'Invalid file type. Allowed: PDF, DOC, DOCX, XLS, XLSX, TXT, CSV, MD, RTF, PPT, PPTX.'


@app.route('/results')
def results_page():
    results = session.get('results')
    if not results:
        return redirect(url_for('home'))
    return render_template(
        'results.html',
        company_name=results['company_name'],
        overview=results['overview'],
        top_keywords=results['top_keywords'],
        theme_counts=results['theme_counts'],
        report_path=results['report_path']
    )


def process_file(filepath: str, file_extension: str) -> Optional[str]:
    """Process different file types and extract their content.

    Args:
        filepath: Path to the file
        file_extension: File extension

    Returns:
        Optional[str]: Extracted content or None if file type is not supported
    """
    try:
        # 1. Handle text-based formats
        if file_extension in ['txt', 'csv', 'md', 'rtf']:
            try:
                with open(filepath, 'r', encoding='utf-8') as f:
                    return f.read()
            except UnicodeDecodeError:
                with open(filepath, 'r', encoding='latin-1') as f:
                    return f.read()

        # 2. Handle PDFs
        elif file_extension == 'pdf':
            from PyPDF2 import PdfReader
            reader = PdfReader(filepath)
            return " ".join([page.extract_text() for page in reader.pages if page.extract_text()])

        # 3. Handle DOCX
        elif file_extension == 'docx':
            import docx
            doc = docx.Document(filepath)
            return "\n".join([para.text for para in doc.paragraphs])

        # 4. Handle DOC
        elif file_extension == 'doc':
            import textract
            return textract.process(filepath).decode('utf-8')

        # 5. Handle XLS/XLSX
        elif file_extension in ['xls', 'xlsx']:
            import pandas as pd
            df = pd.read_excel(filepath, sheet_name=None)
            return "\n".join(df[sheet].to_string(index=False) for sheet in df)

        # 6. Handle PPT/PPTX
        elif file_extension in ['ppt', 'pptx']:
            from pptx import Presentation
            prs = Presentation(filepath)
            content = ""
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        content += shape.text + "\n"
            return content

        return None
    except Exception as e:
        print(f"Error processing file: {e}")
        return None


def contains_email(text: str) -> bool:
    """Detect if the text contains an email address."""
    match = re.search(
        r'\b[A-Za-z0-9._%+-]+@[A-Za-z0-9.-]+\.[A-Z|a-z]{2,}\b', text)
    if match:
        return True
    return False


def contains_phone(text: str) -> bool:
    """Detect if the text contains a phone number (permissive, international and US)."""
    # Accepts numbers with country code, spaces, dashes, parentheses, etc.
    match = re.search(r'(\+?\d[\d\s\-().]{7,}\d)', text)
    if match:
        return True
    return False


# Helper: generate PDF from LLM JSON

def generate_pdf_from_llm_json(report_json, pdf_path, report_name="Business Report"):
    pdf = FPDF()
    pdf.add_page()
    pdf.set_font("Arial", 'B', 16)
    pdf.cell(0, 10, report_name, ln=True, align='C')
    pdf.ln(10)
    pdf.set_font("Arial", '', 12)
    # Add sections
    for section in ["executive_summary", "key_findings", "detailed_analysis", "recommendations", "conclusion"]:
        if section in report_json:
            pdf.set_font("Arial", 'B', 14)
            pdf.cell(0, 10, section.replace('_', ' ').title(), ln=True)
            pdf.set_font("Arial", '', 12)
            pdf.multi_cell(0, 10, report_json[section])
            pdf.ln(5)
    # Add visualizations
    if "visualizations" in report_json:
        for viz in report_json["visualizations"]:
            if viz.get("caption"):
                pdf.set_font("Arial", 'I', 11)
                pdf.multi_cell(0, 8, viz["caption"])
            if viz.get("image_base64"):
                img_data = base64.b64decode(viz["image_base64"])
                img_path = os.path.join(REPORTS_FOLDER, "temp_chart.png")
                with open(img_path, "wb") as img_file:
                    img_file.write(img_data)
                pdf.image(img_path, w=150)
                os.remove(img_path)
            pdf.ln(5)
    pdf.output(pdf_path)

# Helper: email PDF


def email_pdf(pdf_path, report_name):
    recipient_email = os.getenv("RECEIVER_MAIL")
    if not recipient_email:
        return False
    msg = Message(
        subject=f"Trendlyzer Report: {report_name}",
        recipients=[recipient_email],
        body=f"Your Trendlyzer report '{report_name}' is attached."
    )
    with open(pdf_path, 'rb') as fp:
        msg.attach(
            filename=os.path.basename(pdf_path),
            content_type='application/pdf',
            data=fp.read()
        )
    mail.send(msg)
    return True


if __name__ == '__main__':
    port = int(os.environ.get('PORT', 5000))
    app.run(host='0.0.0.0', port=port)
