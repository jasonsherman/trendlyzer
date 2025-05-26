import pdfplumber
import docx
import pandas as pd
from pptx import Presentation
import os


def extract_text_from_pdf(path):
    with pdfplumber.open(path) as pdf:
        return "\n".join(page.extract_text() or "" for page in pdf.pages)


def extract_text_from_docx(path):
    doc = docx.Document(path)
    return "\n".join([p.text for p in doc.paragraphs])


def extract_text_from_txt(path):
    with open(path, encoding="utf-8") as f:
        return f.read()


def extract_text_from_csv(path):
    df = pd.read_csv(path)
    return df.to_string()


def extract_text_from_xlsx(path):
    df = pd.read_excel(path)
    return df.to_string()


def extract_text_from_md(path):
    with open(path, encoding="utf-8") as f:
        return f.read()


def extract_text_from_rtf(path):
    with open(path, encoding="utf-8") as f:
        return f.read()


def extract_text_from_pptx(path):
    prs = Presentation(path)
    text = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text.append(shape.text)
    return "\n".join(text)


def extract_text(path, ext):
    ext = ext.lower()
    if ext == "pdf":
        return extract_text_from_pdf(path)
    elif ext in ["doc", "docx"]:
        return extract_text_from_docx(path)
    elif ext == "txt":
        return extract_text_from_txt(path)
    elif ext == "csv":
        return extract_text_from_csv(path)
    elif ext in ["xls", "xlsx"]:
        return extract_text_from_xlsx(path)
    elif ext == "md":
        return extract_text_from_md(path)
    elif ext == "rtf":
        return extract_text_from_rtf(path)
    elif ext in ["ppt", "pptx"]:
        return extract_text_from_pptx(path)
    else:
        raise ValueError("Unsupported file type")
